from __future__ import annotations

import csv
import hashlib
import io
import json
import xml.etree.ElementTree as ET
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

import httpx
from app.application.ingestion.models import (
    ContentClass,
    DocumentFormat,
    ExtractionStrategy,
    SourceType,
    StructuredDocument,
)
from app.application.ingestion.normalizer import (
    infer_sections_from_text,
    normalize_text,
)


@dataclass
class ProcessorInput:
    filename: str
    content_bytes: bytes | None
    source_url: str | None
    mime_type: str | None
    text_encoding: str | None


def _base_metadata(payload: ProcessorInput) -> dict[str, Any]:
    now = datetime.now(timezone.utc).isoformat()
    raw = payload.content_bytes or b""
    return {
        "filename": payload.filename,
        "mime_type": payload.mime_type,
        "source_url": payload.source_url,
        "file_size": len(raw),
        "checksum": hashlib.sha256(raw).hexdigest() if raw else None,
        "processed_at": now,
    }


def _decode_bytes(raw: bytes, encoding: str | None = None) -> str:
    if not raw:
        return ""
    preferred = [encoding] if encoding else []
    for enc in [*preferred, "utf-8", "utf-16", "ascii", "latin-1"]:
        if not enc:
            continue
        try:
            return raw.decode(enc)
        except Exception:
            continue
    return raw.decode("utf-8", errors="replace")


class TextProcessor:
    def process(self, payload: ProcessorInput) -> StructuredDocument:
        text = _decode_bytes(payload.content_bytes or b"", payload.text_encoding)
        normalized = normalize_text(text)
        metadata = _base_metadata(payload)
        metadata.update(
            {
                "word_count": len(normalized.split()),
                "char_count": len(normalized),
                "line_count": len(normalized.splitlines()),
            }
        )
        return StructuredDocument(
            source_type=SourceType.TEXT,
            raw_text=normalized,
            text=normalized,
            sections=infer_sections_from_text(normalized),
            metadata=metadata,
            confidence_scores={"text_extraction": 0.99},
            confidence_score=0.99,
            extraction_path="text_direct",
            strategy_used=ExtractionStrategy.TEXT_DIRECT,
            format=DocumentFormat.TEXT,
            content_class=ContentClass.TEXT_BASED,
        )


class PdfProcessor:
    def classify_pdf(self, payload: ProcessorInput) -> ContentClass:
        raw = payload.content_bytes or b""
        text_markers = sum(raw.count(marker) for marker in (b"BT", b"/Font", b"Tj"))
        image_markers = sum(raw.count(marker) for marker in (b"/Image", b"/XObject"))
        if text_markers > 0 and image_markers == 0:
            return ContentClass.TEXT_BASED
        if image_markers > 0 and text_markers == 0:
            return ContentClass.IMAGE_BASED
        if image_markers > 0 and text_markers > 0:
            return ContentClass.MIXED
        return ContentClass.TEXT_BASED

    def _extract_text_pdfplumber(
        self, payload: ProcessorInput
    ) -> tuple[str, list[dict[str, Any]], dict[str, Any]]:
        try:
            import pdfplumber
        except Exception:
            return "", [], {"pdfplumber_available": False}

        text_parts: list[str] = []
        tables: list[dict[str, Any]] = []
        page_count = 0

        with pdfplumber.open(io.BytesIO(payload.content_bytes or b"")) as pdf:
            page_count = len(pdf.pages)
            for index, page in enumerate(pdf.pages, start=1):
                page_text = page.extract_text(layout=True) or ""
                if page_text:
                    text_parts.append(f"\n[Page {index}]\n{page_text}")

                extracted_tables = page.extract_tables() or []
                for table in extracted_tables:
                    tables.append({"page": index, "rows": table})

        combined = normalize_text("\n".join(text_parts))
        return (
            combined,
            tables,
            {"pdfplumber_available": True, "page_count": page_count},
        )

    def _extract_with_mineru(
        self,
        payload: ProcessorInput,
    ) -> tuple[str, list[dict[str, Any]], list[dict[str, Any]], dict[str, Any]]:
        try:
            import mineru  # type: ignore
        except Exception:
            return "", [], [], {"mineru_available": False}

        try:
            if hasattr(mineru, "extract"):
                result = mineru.extract(payload.content_bytes or b"")
                text = normalize_text(str(result.get("text") or ""))
                tables = list(result.get("tables") or [])
                images = list(result.get("images") or [])
                return text, tables, images, {"mineru_available": True}
        except Exception as error:
            return "", [], [], {"mineru_available": True, "mineru_error": str(error)}

        return "", [], [], {"mineru_available": True, "mineru_adapter": "missing"}

    def _extract_with_paddleocr(self, payload: ProcessorInput) -> tuple[str, dict[str, Any]]:
        try:
            from paddleocr import PaddleOCR  # type: ignore
            from PIL import Image
        except Exception:
            return "", {"paddleocr_available": False}

        try:
            image = Image.open(io.BytesIO(payload.content_bytes or b""))
            ocr = PaddleOCR(use_angle_cls=True, lang="en")
            result = ocr.ocr(image, cls=True)
            lines: list[str] = []
            for region in result or []:
                for item in region or []:
                    if len(item) >= 2 and item[1]:
                        text_value = item[1][0]
                        score = float(item[1][1]) if len(item[1]) > 1 else 0.0
                        if score >= 0.85:
                            lines.append(str(text_value))
            text = normalize_text("\n".join(lines))
            return text, {"paddleocr_available": True}
        except Exception as error:
            return "", {"paddleocr_available": True, "paddleocr_error": str(error)}

    def _extract_with_gemini_vision(self, payload: ProcessorInput) -> tuple[str, dict[str, Any]]:
        try:
            from google import genai
            from google.genai import types
            from shared.gemini.config import get_gemini_config
        except Exception:
            return "", {"gemini_vision_available": False}

        if not payload.content_bytes:
            return "", {"gemini_vision_available": False, "reason": "empty_content"}

        try:
            import os

            api_key = (os.getenv("GEMINI_API_KEY") or "").strip()
            if not api_key:
                return "", {"gemini_vision_available": False, "reason": "missing_api_key"}

            client = genai.Client(api_key=api_key)
            config = get_gemini_config()
            response = client.models.generate_content(
                model=config.vision_model,
                contents=[
                    types.Part.from_text(
                        "Extract all readable text from this document image. Return plain text only."
                    ),
                    types.Part.from_bytes(data=payload.content_bytes, mime_type=payload.mime_type or "image/png"),
                ],
            )
            text = normalize_text(response.text or "")
            return text, {"gemini_vision_available": True, "gemini_vision_model": config.vision_model}
        except Exception as error:
            return "", {"gemini_vision_available": True, "gemini_vision_error": str(error)}

    def _extract_image_pdf_fallback(
        self, payload: ProcessorInput, content_class: ContentClass
    ) -> StructuredDocument:
        mineru_text, mineru_tables, mineru_images, mineru_meta = self._extract_with_mineru(payload)
        paddle_text, paddle_meta = self._extract_with_paddleocr(payload)

        tesseract_text = ""
        tesseract_meta: dict[str, Any] = {}
        try:
            import pytesseract
            from PIL import Image

            image = Image.open(io.BytesIO(payload.content_bytes or b""))
            tesseract_text = normalize_text(pytesseract.image_to_string(image))
            tesseract_meta = {"tesseract_available": True}
        except Exception as error:
            tesseract_meta = {"tesseract_available": False, "tesseract_error": str(error)}

        gemini_text, gemini_meta = self._extract_with_gemini_vision(payload)

        text = mineru_text or paddle_text or tesseract_text or gemini_text
        tables = mineru_tables
        images = mineru_images

        metadata = _base_metadata(payload)
        metadata.update(
            {
                "ocr_chain": ["mineru", "paddleocr", "tesseract", "gemini_vision"],
                **mineru_meta,
                **paddle_meta,
                **tesseract_meta,
                **gemini_meta,
            }
        )
        confidence = 0.9 if text else 0.2
        return StructuredDocument(
            source_type=(
                SourceType.PDF_SCANNED
                if content_class == ContentClass.IMAGE_BASED
                else SourceType.PDF_MIXED
            ),
            raw_text=text,
            text=text,
            sections=infer_sections_from_text(text),
            tables=tables,
            images=images,
            metadata=metadata,
            confidence_scores={"text_extraction": confidence},
            confidence_score=confidence,
            extraction_path=(
                "mineru->paddleocr->tesseract->gemini_vision"
            ),
            strategy_used=(
                ExtractionStrategy.PDF_SCANNED_HIGH_ACCURACY
                if content_class == ContentClass.IMAGE_BASED
                else ExtractionStrategy.PDF_MIXED_HYBRID
            ),
            format=DocumentFormat.PDF,
            content_class=content_class,
        )

    def process(self, payload: ProcessorInput) -> StructuredDocument:
        content_class = self.classify_pdf(payload)

        if content_class == ContentClass.TEXT_BASED:
            try:
                text, tables, extra_meta = self._extract_text_pdfplumber(payload)
            except Exception:
                return self._extract_image_pdf_fallback(payload, ContentClass.MIXED)
            metadata = _base_metadata(payload)
            metadata.update(extra_meta)
            metadata.update(
                {
                    "word_count": len(text.split()),
                    "char_count": len(text),
                }
            )
            return StructuredDocument(
                source_type=SourceType.PDF_TEXT,
                raw_text=text,
                text=text,
                sections=infer_sections_from_text(text),
                tables=tables,
                metadata=metadata,
                confidence_scores={"text_extraction": 0.95, "table_extraction": 0.7},
                confidence_score=0.95,
                extraction_path="pdfplumber",
                strategy_used=ExtractionStrategy.PDF_TEXT_FAST,
                format=DocumentFormat.PDF,
                content_class=content_class,
            )

        return self._extract_image_pdf_fallback(payload, content_class)


class ImageProcessor:
    def process(self, payload: ProcessorInput) -> StructuredDocument:
        metadata = _base_metadata(payload)
        metadata["ocr_chain"] = ["paddleocr", "tesseract", "gemini_vision"]

        text = ""
        confidence = 0.2

        try:
            import pytesseract
            from PIL import Image

            image = Image.open(io.BytesIO(payload.content_bytes or b""))
            text = normalize_text(pytesseract.image_to_string(image))
            confidence = 0.7 if text else 0.3
            metadata["ocr_engine"] = "tesseract"
        except Exception:
            metadata["ocr_engine"] = "fallback_chain_unavailable"

        return StructuredDocument(
            source_type=SourceType.IMAGE,
            raw_text=text,
            text=text,
            sections=infer_sections_from_text(text),
            metadata=metadata,
            confidence_scores={"ocr_confidence": confidence},
            confidence_score=confidence,
            extraction_path="paddleocr->tesseract->gemini_vision",
            strategy_used=ExtractionStrategy.IMAGE_OCR_CHAIN,
            format=DocumentFormat.IMAGE,
            content_class=ContentClass.IMAGE_BASED,
        )


class OfficeProcessor:
    def _process_docx(self, payload: ProcessorInput) -> tuple[str, dict[str, Any]]:
        try:
            from docx import Document
        except Exception:
            return "", {"docx_parser": "unavailable"}

        document = Document(io.BytesIO(payload.content_bytes or b""))
        paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
        table_count = len(document.tables)
        return normalize_text("\n\n".join(paragraphs)), {"table_count": table_count}

    def _process_xlsx(self, payload: ProcessorInput) -> tuple[str, dict[str, Any]]:
        try:
            from openpyxl import load_workbook
        except Exception:
            return "", {"xlsx_parser": "unavailable"}

        workbook = load_workbook(
            io.BytesIO(payload.content_bytes or b""), data_only=False
        )
        markdown_parts: list[str] = []
        for sheet in workbook.worksheets:
            markdown_parts.append(f"## Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                values = ["" if cell is None else str(cell) for cell in row]
                markdown_parts.append(" | ".join(values))
        return normalize_text("\n".join(markdown_parts)), {
            "sheet_count": len(workbook.worksheets)
        }

    def _process_pptx(self, payload: ProcessorInput) -> tuple[str, dict[str, Any]]:
        try:
            from pptx import Presentation
        except Exception:
            return "", {"pptx_parser": "unavailable"}

        presentation = Presentation(io.BytesIO(payload.content_bytes or b""))
        lines: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            lines.append(f"## Slide {index}")
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    lines.append(text.strip())
        return normalize_text("\n".join(lines)), {
            "slide_count": len(presentation.slides)
        }

    def process(self, payload: ProcessorInput) -> StructuredDocument:
        extension = Path(payload.filename).suffix.lower()
        text = ""
        details: dict[str, Any] = {}

        if extension in {".docx", ".doc"}:
            text, details = self._process_docx(payload)
        elif extension in {".xlsx", ".xls"}:
            text, details = self._process_xlsx(payload)
        elif extension in {".pptx", ".ppt"}:
            text, details = self._process_pptx(payload)

        metadata = _base_metadata(payload)
        metadata.update(details)
        metadata["office_extension"] = extension

        return StructuredDocument(
            source_type=(
                SourceType.DOCX
                if extension in {".docx", ".doc"}
                else SourceType.XLSX
                if extension in {".xlsx", ".xls"}
                else SourceType.PPTX
            ),
            raw_text=text,
            text=text,
            sections=infer_sections_from_text(text),
            metadata=metadata,
            confidence_scores={"text_extraction": 0.9 if text else 0.2},
            confidence_score=0.9 if text else 0.2,
            extraction_path="office_native",
            strategy_used=ExtractionStrategy.OFFICE_NATIVE,
            format=DocumentFormat.OFFICE,
            content_class=ContentClass.TEXT_BASED,
        )


class WebProcessor:
    async def _fetch_html(self, url: str) -> str:
        timeout = httpx.Timeout(10.0)
        async with httpx.AsyncClient(timeout=timeout, follow_redirects=True) as client:
            response = await client.get(url)
            response.raise_for_status()
            return response.text

    def _extract_static(self, html: str) -> tuple[str, dict[str, Any]]:
        try:
            import trafilatura

            extracted = trafilatura.extract(html, output_format="txt")
            if extracted:
                return normalize_text(extracted), {"extractor": "trafilatura"}
        except Exception:
            pass

        try:
            from bs4 import BeautifulSoup
            from readability import Document

            summary_html = Document(html).summary()
            text = BeautifulSoup(summary_html, "lxml").get_text("\n")
            return normalize_text(text), {"extractor": "readability"}
        except Exception:
            pass

        from bs4 import BeautifulSoup

        text = BeautifulSoup(html, "lxml").get_text("\n")
        return normalize_text(text), {"extractor": "beautifulsoup"}

    async def _render_dynamic(self, url: str) -> tuple[str, dict[str, Any]]:
        try:
            from playwright.async_api import async_playwright
        except Exception:
            return "", {"playwright_available": False}

        try:
            async with async_playwright() as playwright:
                browser = await playwright.chromium.launch(headless=True)
                context = await browser.new_context()
                page = await context.new_page()
                async def handle_route(route):
                    if route.request.resource_type in {"image", "font", "media"}:
                        await route.abort()
                    else:
                        await route.continue_()

                await page.route("**/*", handle_route)
                await page.goto(url, wait_until="domcontentloaded", timeout=30000)
                await page.wait_for_load_state("networkidle", timeout=30000)
                html = await page.content()
                await context.close()
                await browser.close()
        except Exception as error:
            return "", {"playwright_available": True, "playwright_error": str(error)}

        text, details = self._extract_static(html)
        merged = {"playwright_available": True, "extractor": f"playwright->{details.get('extractor', 'static')}"}
        merged.update(details)
        return text, merged

    @staticmethod
    def _is_dynamic_candidate(html: str) -> bool:
        sample = html.lower()
        signals = ["id=\"root\"", "id='root'", "__next", "data-reactroot", "ng-app", "vue"]
        text_density = len("".join(ch for ch in sample[:4000] if ch.isalpha())) / max(1, len(sample[:4000]))
        return any(signal in sample for signal in signals) or text_density < 0.2

    async def process(self, payload: ProcessorInput) -> StructuredDocument:
        if not payload.source_url:
            return StructuredDocument(
                source_type=SourceType.WEB_STATIC,
                raw_text="",
                text="",
                metadata={"error": "missing source_url"},
                confidence_scores={"text_extraction": 0.0},
                confidence_score=0.0,
                extraction_path="web_missing_url",
                strategy_used=ExtractionStrategy.WEB_STATIC_DYNAMIC,
                format=DocumentFormat.WEB,
            )

        try:
            html = await self._fetch_html(payload.source_url)
            if self._is_dynamic_candidate(html):
                dynamic_text, dynamic_details = await self._render_dynamic(payload.source_url)
                if dynamic_text:
                    text, details = dynamic_text, dynamic_details
                else:
                    static_text, static_details = self._extract_static(html)
                    text = static_text
                    details = {**dynamic_details, **static_details}
            else:
                text, details = self._extract_static(html)
        except Exception as error:
            text = ""
            details = {"extractor": "network_fallback", "error": str(error)}
        metadata = _base_metadata(payload)
        metadata.update(details)
        metadata["word_count"] = len(text.split())

        source_type = (
            SourceType.WEB_DYNAMIC
            if str(details.get("extractor", "")).startswith("playwright")
            else SourceType.WEB_STATIC
        )

        return StructuredDocument(
            source_type=source_type,
            raw_text=text,
            text=text,
            sections=infer_sections_from_text(text),
            metadata=metadata,
            confidence_scores={"text_extraction": 0.9 if text else 0.2},
            confidence_score=0.9 if text else 0.2,
            extraction_path=str(details.get("extractor", "web_unknown")),
            strategy_used=ExtractionStrategy.WEB_STATIC_DYNAMIC,
            format=DocumentFormat.WEB,
            content_class=ContentClass.TEXT_BASED,
        )


class FallbackProcessor:
    def process(self, payload: ProcessorInput) -> StructuredDocument:
        text = _decode_bytes(payload.content_bytes or b"", payload.text_encoding)
        normalized = normalize_text(text)
        metadata = _base_metadata(payload)
        return StructuredDocument(
            source_type=SourceType.TEXT,
            raw_text=normalized,
            text=normalized,
            sections=infer_sections_from_text(normalized),
            metadata=metadata,
            confidence_scores={"text_extraction": 0.4},
            confidence_score=0.4,
            extraction_path="fallback_generic",
            strategy_used=ExtractionStrategy.FALLBACK_GENERIC,
            format=DocumentFormat.UNKNOWN,
        )


def parse_text_subtype(text: str, subtype: str | None) -> dict[str, Any]:
    if not subtype:
        return {}
    if subtype == "json":
        try:
            data = json.loads(text)
            return {"json_keys": list(data.keys()) if isinstance(data, dict) else None}
        except Exception:
            return {}
    if subtype == "xml":
        try:
            root = ET.fromstring(text)
            return {"xml_root_tag": root.tag}
        except Exception:
            return {}
    if subtype == "csv":
        try:
            reader = csv.reader(io.StringIO(text))
            first = next(reader, [])
            return {"csv_header": first}
        except Exception:
            return {}
    return {}
